from pptx import Presentation
from pptx.util import Pt

def create_presentation():
    # Load the Red Hat template you downloaded
    try:
        prs = Presentation("redhat_template.pptx")
    except Exception as e:
        print(f"Error loading template. Ensure 'redhat_template.pptx' is in the same folder. Error: {e}")
        return

    # Slide Data based on the outline
    slides_data = [
        {
            "type": "title",
            "title": "Bring Your Own Knowledge to OpenShift Lightspeed",
            "content": "Tailoring AI to Your Organization's Unique Context"
        },
        {
            "type": "content",
            "title": "The AI Knowledge Gap",
            "content": "• Organizations have deep, proprietary knowledge (runbooks, configs, compliance rules).\n• Public Large Language Models (LLMs) don't have access to this data.\n• Result: AI provides generic guidance that may not align with your internal policies."
        },
        {
            "type": "content",
            "title": "Bridge the Gap: Bring Your Own Knowledge",
            "content": "• Augment Lightspeed's intelligence with your private documentation.\n• Transform AI from a generic tool to a context-aware partner.\n• Deliver tailored, policy-compliant answers."
        },
        {
            "type": "content",
            "title": "Intelligently Tailored for Any Industry",
            "content": "• Finance: Ingest internal security policies and compliance checklists.\n• Telecom: Support bespoke network configurations.\n• Government: Adhere to unique procedural requirements and governance."
        },
        {
            "type": "content",
            "title": "How It Works (Steps 1 & 2)",
            "content": "1. Gather your documentation into Markdown (.md) files.\n2. Use the lightspeed-rag-tool container to compile them.\n3. Output a custom container image (byok-image.tar).\n4. Push the image to your OpenShift registry."
        },
        {
            "type": "content",
            "title": "How It Works (Step 3)",
            "content": "Configure OpenShift Lightspeed (OLSConfig):\n\nols:\n  defaultModel: gpt-4o\n  rag:\n    - image: image-registry.../acme-byok:latest\n      indexID: vector_db_index\n      indexPath: /rag/vector_db"
        },
        {
            "type": "content",
            "title": "Demo Time: Let's See It In Action",
            "content": "• Goal: Asking OpenShift Lightspeed an internal policy question.\n• Watch as the AI references our newly ingested, private documentation to provide a tailored answer."
        },
        {
            "type": "content",
            "title": "Wrap-Up & Next Steps",
            "content": "• BYO Knowledge is currently in Technology Preview.\n• Expect rapid improvements as the feature matures toward General Availability.\n• Try OpenShift Lightspeed in your environment today!\n• Visit: redhat.com/openshift"
        }
    ]

    # Add slides to the presentation
    for slide_info in slides_data:
        if slide_info["type"] == "title":
            # Usually, layout 0 is the Title Slide layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_info["title"]
            subtitle.text = slide_info["content"]
            
        elif slide_info["type"] == "content":
            # Usually, layout 1 is the Title and Content layout
            # If the Red Hat template uses a different index for standard content slides, change the '1' below.
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_info["title"]
            body.text = slide_info["content"]

    # Save the new presentation
    output_filename = "OpenShift_Lightspeed_BYOK_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}!")

if __name__ == "__main__":
    create_presentation()
